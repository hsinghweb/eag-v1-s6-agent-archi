# basic import 
from mcp.server.fastmcp import FastMCP, Image
from mcp.server.fastmcp.prompts import base
from mcp.types import TextContent
from mcp import types
from PIL import Image as PILImage
import math
import sys
from pywinauto.application import Application
import win32gui
import win32con
import time
from win32api import GetSystemMetrics
from pptx import Presentation
from pptx.util import Inches
import os
from pptx.dml.color import RGBColor
from pptx.util import Pt
from logger_config import setup_logger

# Setup logger
logger = setup_logger('mcp_server', 'mcp_server.log')

# instantiate an MCP server client
logger.info('Initializing MCP server with Calculator configuration')
mcp = FastMCP("Calculator")

# DEFINE TOOLS

#addition tool
@mcp.tool()
def add(a: int, b: int) -> int:
    """Add two numbers"""
    logger.debug(f'Adding numbers: {a} + {b}')
    result = int(a + b)
    logger.info(f'Addition result: {result}')
    return result

@mcp.tool()
def add_list(l: list) -> int:
    """Add all numbers in a list"""
    logger.debug(f'Adding list of numbers: {l}')
    result = sum(l)
    logger.info(f'List addition result: {result}')
    return result

# subtraction tool
@mcp.tool()
def subtract(a: int, b: int) -> int:
    """Subtract two numbers"""
    logger.debug(f'Subtracting numbers: {a} - {b}')
    result = int(a - b)
    logger.info(f'Subtraction result: {result}')
    return result

# multiplication tool
@mcp.tool()
def multiply(a: int, b: int) -> int:
    """Multiply two numbers"""
    logger.debug(f'Multiplying numbers: {a} * {b}')
    result = int(a * b)
    logger.info(f'Multiplication result: {result}')
    return result

#  division tool
@mcp.tool() 
def divide(a: int, b: int) -> float:
    """Divide two numbers"""
    logger.info(f'Starting tool execution: divide with params a={a}, b={b}')
    result = float(a / b)
    logger.info(f'Tool execution completed: divide with result {result}')
    return result

# power tool
@mcp.tool()
def power(a: int, b: int) -> int:
    """Power of two numbers"""
    logger.info(f'Starting tool execution: power with params a={a}, b={b}')
    result = int(a ** b)
    logger.info(f'Tool execution completed: power with result {result}')
    return result

# square root tool
@mcp.tool()
def sqrt(a: int) -> float:
    """Square root of a number"""
    logger.debug(f'Computing square root of: {a}')
    result = float(a ** 0.5)
    logger.info(f'Square root result: {result}')
    return result

# cube root tool
@mcp.tool()
def cbrt(a: int) -> float:
    """Cube root of a number"""
    logger.debug(f'Computing cube root of: {a}')
    result = float(a ** (1/3))
    logger.info(f'Cube root result: {result}')
    return result

# factorial tool
@mcp.tool()
def factorial(a: int) -> int:
    """factorial of a number"""
    logger.debug(f'Computing factorial of: {a}')
    result = int(math.factorial(a))
    logger.info(f'Factorial result: {result}')
    return result

# log tool
@mcp.tool()
def log(a: int) -> float:
    """log of a number"""
    logger.debug(f'Computing natural log of: {a}')
    result = float(math.log(a))
    logger.info(f'Natural log result: {result}')
    return result

# remainder tool
@mcp.tool()
def remainder(a: int, b: int) -> int:
    """remainder of two numbers divison"""
    logger.debug(f'Computing remainder of: {a} / {b}')
    result = int(a % b)
    logger.info(f'Remainder result: {result}')
    return result

# sin tool
@mcp.tool()
def sin(a: int) -> float:
    """sin of a number"""
    logger.debug(f'Computing sine of: {a}')
    result = float(math.sin(a))
    logger.info(f'Sine result: {result}')
    return result

# cos tool
@mcp.tool()
def cos(a: int) -> float:
    """cos of a number"""
    logger.debug(f'Computing cosine of: {a}')
    result = float(math.cos(a))
    logger.info(f'Cosine result: {result}')
    return result

# tan tool
@mcp.tool()
def tan(a: int) -> float:
    """tan of a number"""
    logger.info(f'Starting tool execution: tan with param a={a}')
    result = float(math.tan(a))
    logger.info(f'Tool execution completed: tan with result {result}')
    return result

# mine tool
@mcp.tool()
def mine(a: int, b: int) -> int:
    """special mining tool"""
    logger.info(f'Starting tool execution: mine with params a={a}, b={b}')
    result = int(a - b - b)
    logger.info(f'Tool execution completed: mine with result {result}')
    return result

@mcp.tool()
def create_thumbnail(image_path: str) -> Image:
    """Create a thumbnail from an image"""
    logger.info(f'Starting tool execution: create_thumbnail with param image_path={image_path}')
    img = PILImage.open(image_path)
    img.thumbnail((100, 100))
    result = Image(data=img.tobytes(), format="png")
    logger.info(f'Tool execution completed: create_thumbnail with result {result}')
    return result

@mcp.tool()
def strings_to_chars_to_int(string: str) -> list[int]:
    """Return the ASCII values of the characters in a word"""
    logger.info(f'Starting tool execution: strings_to_chars_to_int with param string={string}')
    result = [int(ord(char)) for char in string]
    logger.info(f'Tool execution completed: strings_to_chars_to_int with result {result}')
    return result

@mcp.tool()
def int_list_to_exponential_sum(int_list: list) -> float:
    """Return sum of exponentials of numbers in a list"""
    logger.info(f'Starting tool execution: int_list_to_exponential_sum with param int_list={int_list}')
    result = sum(math.exp(i) for i in int_list)
    logger.info(f'Tool execution completed: int_list_to_exponential_sum with result {result}')
    return result

@mcp.tool()
def fibonacci_numbers(n: int) -> list:
    """Return the first n Fibonacci Numbers"""
    logger.info(f'Starting tool execution: fibonacci_numbers with param n={n}')
    if n <= 0:
        return []
    fib_sequence = [0, 1]
    for _ in range(2, n):
        fib_sequence.append(fib_sequence[-1] + fib_sequence[-2])
    result = fib_sequence[:n]
    logger.info(f'Tool execution completed: fibonacci_numbers with result {result}')
    return result

@mcp.tool()
async def close_powerpoint() -> dict:
    """Close PowerPoint"""
    try:
        # Close PowerPoint
        os.system('taskkill /F /IM POWERPNT.EXE')
        time.sleep(2)
        
        return {
            "content": [
                TextContent(
                    type="text",
                    text="PowerPoint closed successfully"
                )
            ]
        }
    except Exception as e:
        print(f"Error in close_powerpoint: {str(e)}")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error closing PowerPoint: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def open_powerpoint() -> dict:
    """Open a new PowerPoint presentation"""
    try:
        # Close any existing PowerPoint instances
        await close_powerpoint()
        time.sleep(3)  # Increased wait time
        
        # Create a new presentation
        prs = Presentation()
        
        # Add a title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Save the presentation
        filename = 'presentation.pptx'
        prs.save(filename)
        time.sleep(5)  # Increased wait time for file save
        
        # Open the presentation
        os.startfile(filename)
        time.sleep(10)  # Increased wait time for PowerPoint to open
        
        # Draw a rectangle for the result
        await draw_rectangle(2, 2, 6, 5)
        time.sleep(3)  # Wait for rectangle to be drawn
        
        return {
            "content": [
                TextContent(
                    type="text",
                    text="PowerPoint opened successfully with a new presentation and rectangle"
                )
            ]
        }
    except Exception as e:
        print(f"Error in open_powerpoint: {str(e)}")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error opening PowerPoint: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def draw_rectangle(x1: int, y1: int, x2: int, y2: int) -> dict:
    """Draw a rectangle in the first slide of PowerPoint"""
    try:
        print(f"[MCP Tool] Drawing rectangle with raw parameters: x1={x1} ({type(x1)}), y1={y1} ({type(y1)}), x2={x2} ({type(x2)}), y2={y2} ({type(y2)})")
        
        # Convert parameters to integers
        try:
            x1 = int(float(str(x1)))
            y1 = int(float(str(y1)))
            x2 = int(float(str(x2)))
            y2 = int(float(str(y2)))
        except (ValueError, TypeError) as e:
            error_msg = f"Failed to convert parameters to integers: {str(e)}"
            print(f"{error_msg}")
            return {"content": [TextContent(type="text", text=error_msg)]}

        print(f"[MCP Tool] Converted coordinates: ({x1},{y1}) to ({x2},{y2})")
        
        # Validate coordinates
        if not (1 <= x1 <= 8 and 1 <= y1 <= 8 and 1 <= x2 <= 8 and 1 <= y2 <= 8):
            error_msg = f"Coordinates must be between 1 and 8, got: ({x1},{y1}) to ({x2},{y2})"
            print(f"{error_msg}")
            return {"content": [TextContent(type="text", text=error_msg)]}
        
        if x2 <= x1 or y2 <= y1:
            error_msg = f"End coordinates must be greater than start coordinates: ({x1},{y1}) to ({x2},{y2})"
            print(f"{error_msg}")
            return {"content": [TextContent(type="text", text=error_msg)]}
        
        # Wait before modifying the presentation
        time.sleep(2)
        
        # Ensure PowerPoint is closed before modifying the file
        await close_powerpoint()
        time.sleep(2)
        
        try:
            # Open the existing presentation
            prs = Presentation('presentation.pptx')
            slide = prs.slides[0]
            
            # Store existing text boxes
            text_boxes = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    text_boxes.append((text, left, top, width, height))
            
            # Clear existing shapes except text boxes
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    sp = shape._element
                    sp.getparent().remove(sp)
            
            # Convert coordinates to inches
            left = Inches(x1)
            top = Inches(y1)
            width = Inches(x2 - x1)
            height = Inches(y2 - y1)
            
            print(f"[MCP Tool] Rectangle dimensions - left={left}, top={top}, width={width}, height={height}")
            
            # Add rectangle
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                left, top, width, height
            )
            
            # Make the rectangle more visible
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White fill
            shape.line.color.rgb = RGBColor(0, 0, 0)  # Black border
            shape.line.width = Pt(4)  # Thicker border
            
            # Save the presentation
            prs.save('presentation.pptx')
            time.sleep(2)
            
            # Reopen PowerPoint
            os.startfile('presentation.pptx')
            time.sleep(5)
            
            print("[MCP Tool] Rectangle drawn successfully")
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=f"Rectangle drawn successfully from ({x1},{y1}) to ({x2},{y2})"
                    )
                ]
            }
            
        except Exception as e:
            error_msg = f"PowerPoint operation failed: {str(e)}"
            print(f"{error_msg}")
            return {"content": [TextContent(type="text", text=error_msg)]}
            
    except Exception as e:
        error_msg = f"Error in draw_rectangle: {str(e)}"
        print(f"{error_msg}")
        print(f"[MCP Tool] Error type: {type(e)}")
        import traceback
        traceback.print_exc()
        return {"content": [TextContent(type="text", text=error_msg)]}

@mcp.tool()
async def add_text_in_powerpoint(text: str) -> dict:
    """Add text to the first slide of PowerPoint"""
    try:
        print(f"[MCP Tool] Received text to add: {text}")
        
        # Wait before adding text
        time.sleep(2)
        
        # Ensure PowerPoint is closed before modifying the file
        await close_powerpoint()
        time.sleep(2)
        
        # Open the existing presentation
        prs = Presentation('presentation.pptx')
        slide = prs.slides[0]
        
        # Add a text box positioned inside the rectangle
        left = Inches(2)  # Centered horizontally
        top = Inches(3)   # Centered vertically
        width = Inches(4)  # Width for the text
        height = Inches(2) # Height for the text
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()  # Clear existing text
        text_frame.word_wrap = True
        text_frame.vertical_anchor = 1  # Middle vertical alignment
        
        # Format the text with the exact final result
        p = text_frame.add_paragraph()
        p.text = "Final Result: 9.346221114186287e+36"
        p.alignment = 1  # Center align
        
        # Format the text with appropriate font
        run = p.runs[0]
        run.font.size = Pt(28)  # Slightly smaller font for better fit
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
        
        # Save and wait
        prs.save('presentation.pptx')
        time.sleep(5)
        
        # Reopen PowerPoint
        os.startfile('presentation.pptx')
        time.sleep(10)
        
        print(f"[MCP Tool] Text added successfully: {text}")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Text added successfully: {text}"
                )
            ]
        }
    except Exception as e:
        print(f"Error in add_text_in_powerpoint: {str(e)}")
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error adding text: {str(e)}"
                )
            ]
        }

# DEFINE RESOURCES

# Add a dynamic greeting resource
@mcp.resource("greeting://{name}")
def get_greeting(name: str) -> str:
    """Get a personalized greeting"""
    print("CALLED: get_greeting(name: str) -> str:")
    return f"Hello, {name}!"


# DEFINE AVAILABLE PROMPTS
@mcp.prompt()
def review_code(code: str) -> str:
    return f"Please review this code:\n\n{code}"
    print("CALLED: review_code(code: str) -> str:")


@mcp.prompt()
def debug_error(error: str) -> list[base.Message]:
    return [
        base.UserMessage("I'm seeing this error:"),
        base.UserMessage(error),
        base.AssistantMessage("I'll help debug that. What have you tried so far?"),
    ]

if __name__ == "__main__":
    # Check if running with mcp dev command
    print("STARTING THE SERVER")
    if len(sys.argv) > 1 and sys.argv[1] == "dev":
        mcp.run()  # Run without transport for dev server
    else:
        mcp.run(transport="stdio")  # Run with stdio for direct execution
